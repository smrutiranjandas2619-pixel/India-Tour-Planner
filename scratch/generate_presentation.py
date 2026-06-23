import os
import sys

def check_dependencies():
    try:
        import pptx
    except ImportError:
        print("[!] python-pptx is not installed. Installing it now...")
        import subprocess
        try:
            subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx"], check=True)
            print("[+] python-pptx installed successfully!\n")
        except Exception as e:
            print(f"[-] Failed to install python-pptx: {e}")
            print("Please run: pip install python-pptx")
            sys.exit(1)

check_dependencies()

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Define color scheme (Deep Slate Blue, Teal/Cyan, Light Gray)
    COLOR_BG = RGBColor(18, 24, 38)        # #121826 (Dark background)
    COLOR_TITLE = RGBColor(0, 212, 255)     # #00D4FF (Cyan accent)
    COLOR_TEXT = RGBColor(220, 225, 235)    # #DCE1EB (Light gray text)
    COLOR_MUTED = RGBColor(150, 160, 175)   # #96A0AF (Muted gray)
    
    # Slide Dimensions (16:9 Widescreen)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Helper function to style text frame
    def style_text_frame(tf):
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)

    # Helper function to set slide background color
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    # Helper function to add a title to a slide
    def add_slide_header(slide, title_text, category_text="INDIA TOUR PLANNER"):
        # Add Category Muted Text
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        cat_tf = cat_box.text_frame
        style_text_frame(cat_tf)
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = category_text.upper()
        cat_p.font.name = "Arial"
        cat_p.font.size = Pt(11)
        cat_p.font.bold = True
        cat_p.font.color.rgb = COLOR_MUTED
        
        # Add Main Title Text
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        title_tf = title_box.text_frame
        style_text_frame(title_tf)
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.name = "Arial"
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = COLOR_TITLE

    # Blank layout for full customization
    slide_layout = prs.slide_layouts[6]

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Cover Slide)
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    
    accent_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf = accent_box.text_frame
    style_text_frame(tf)
    
    p_pre = tf.paragraphs[0]
    p_pre.text = "AI ITINERARY & ADVISORY PLATFORM"
    p_pre.font.name = "Arial"
    p_pre.font.size = Pt(14)
    p_pre.font.bold = True
    p_pre.font.color.rgb = COLOR_TITLE
    p_pre.space_after = Pt(10)
    
    p_title = tf.add_paragraph()
    p_title.text = "India Tour Planner"
    p_title.font.name = "Arial"
    p_title.font.size = Pt(56)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255)
    p_title.space_after = Pt(15)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "A Smart, RAG-Powered Multi-Tier Travel Ecosystem"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = COLOR_TEXT
    
    # Footer info
    footer_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.33), Inches(0.8))
    ftf = footer_box.text_frame
    style_text_frame(ftf)
    fp = ftf.paragraphs[0]
    fp.text = "FastAPI Backend • React (Vite) Frontend • Gemini & Llama 3"
    fp.font.name = "Arial"
    fp.font.size = Pt(12)
    fp.font.color.rgb = COLOR_MUTED

    # -------------------------------------------------------------
    # SLIDE 2: Introduction
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "The Challenge of Travel Planning", "01. Introduction")
    
    # Two Columns
    col1_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf1 = col1_box.text_frame
    style_text_frame(tf1)
    
    p1 = tf1.paragraphs[0]
    p1.text = "The Problem"
    p1.font.name = "Arial"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TITLE
    p1.space_after = Pt(10)
    
    bullets1 = [
        "Planning a trip to India is exciting, but organizing everything takes hours of tedious manual research.",
        "Travelers have to visit separate websites to find sights, check local weather, look up hotel prices, and search safety warnings.",
        "Static templates do not adjust to different budgets or react to changing local climate warnings."
    ]
    for b in bullets1:
        bp = tf1.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)
        
    col2_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf2 = col2_box.text_frame
    style_text_frame(tf2)
    
    p2 = tf2.paragraphs[0]
    p2.text = "Our Simple Solution"
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TITLE
    p2.space_after = Pt(10)
    
    bullets2 = [
        "We built a single web dashboard that coordinates all planning in seconds.",
        "Using Artificial Intelligence, the system merges guides and scrapers into one view.",
        "Generates customized itineraries, plots interactive maps, and calculates realistic budgets automatically."
    ]
    for b in bullets2:
        bp = tf2.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 3: Overview
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "How the Planner Works", "02. Overview")
    
    # 3 Columns
    col_width = Inches(3.6)
    col_gap = Inches(0.4)
    
    col_data = [
        ("1. User Inputs", "Set Your Preferences", [
            "Input destination, trip duration, budget tier (Low, Medium, Premium), and start city.",
            "Helps target exact parameters tailored to individual travel expectations."
        ]),
        ("2. AI Engine", "Smart Synthesis", [
            "AI retrieves context from pre-loaded local vector travel guides.",
            "Combines coordinates and generates structured day-by-day itineraries."
        ]),
        ("3. Visual Deck", "Explore & Tune", [
            "Displays paths on an interactive Leaflet map.",
            "Displays dynamic budget breakdowns.",
            "Includes a chatbot for instant follow-up questions."
        ])
    ]
    
    for i, (title, subtitle, items) in enumerate(col_data):
        col_left = Inches(0.8) + i * (col_width + col_gap)
        c_box = slide.shapes.add_textbox(col_left, Inches(1.8), col_width, Inches(4.5))
        ctf = c_box.text_frame
        style_text_frame(ctf)
        
        tp = ctf.paragraphs[0]
        tp.text = title
        tp.font.name = "Arial"
        tp.font.size = Pt(18)
        tp.font.bold = True
        tp.font.color.rgb = COLOR_TITLE
        tp.space_after = Pt(4)
        
        stp = ctf.add_paragraph()
        stp.text = subtitle
        stp.font.name = "Arial"
        stp.font.size = Pt(13)
        stp.font.italic = True
        stp.font.color.rgb = COLOR_MUTED
        stp.space_after = Pt(12)
        
        for item in items:
            itp = ctf.add_paragraph()
            itp.text = "▪ " + item
            itp.font.name = "Arial"
            itp.font.size = Pt(13)
            itp.font.color.rgb = COLOR_TEXT
            itp.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 4: Key Features
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "What Can the App Do?", "03. Key Features")
    
    # 2 Column layout
    col1_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf1 = col1_box.text_frame
    style_text_frame(tf1)
    
    p1 = tf1.paragraphs[0]
    p1.text = "Smart Planning Features"
    p1.font.name = "Arial"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TITLE
    p1.space_after = Pt(10)
    
    bullets1 = [
        "AI Itinerary Architect: Custom day-wise plans matching user-specified location coordinates.",
        "Smart Expense Estimator: Dynamic budget calculator split by dining, hotels, and transport.",
        "Live Weather & Safety Scraper: Pulls alerts to warn you of storms or weather dangers."
    ]
    for b in bullets1:
        bp = tf1.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)
        
    col2_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf2 = col2_box.text_frame
    style_text_frame(tf2)
    
    p2 = tf2.paragraphs[0]
    p2.text = "User Conveniences"
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TITLE
    p2.space_after = Pt(10)
    
    bullets2 = [
        "Firebase Phone OTP: SMS-based passwordless log in for quick, secure onboarding.",
        "Dossier History Portal: Save, reload, review, and delete previous travel plans directly.",
        "Interactive Chatbot: Embedded assistant answers queries like 'What clothes should I pack?'"
    ]
    for b in bullets2:
        bp = tf2.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 5: Technology Stack
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "The Tools We Used", "04. Technology Stack")
    
    # 3 Columns
    col_width = Inches(3.6)
    col_gap = Inches(0.4)
    
    col_data = [
        ("Frontend Client", "What the User Sees", [
            "React.js Core SPA",
            "Vite compiler framework",
            "Tailwind CSS Layouts",
            "Vanilla CSS variables",
            "Leaflet JS Maps"
        ]),
        ("Backend Services", "The Engine Layer", [
            "FastAPI async routes",
            "Uvicorn ASGI launcher",
            "SQLite database tables",
            "JWT token security codes",
            "Python weather scrapers"
        ]),
        ("AI Core", "The Brain Layer", [
            "Gemini 1.5 (Itinerary details)",
            "Llama 3.3 (High-speed chatbot)",
            "SimpleVectorStore (Cosine DB)",
            "NumPy similarity math",
            "Preloaded Local Guide texts"
        ])
    ]
    
    for i, (title, subtitle, items) in enumerate(col_data):
        col_left = Inches(0.8) + i * (col_width + col_gap)
        c_box = slide.shapes.add_textbox(col_left, Inches(1.8), col_width, Inches(4.5))
        ctf = c_box.text_frame
        style_text_frame(ctf)
        
        tp = ctf.paragraphs[0]
        tp.text = title
        tp.font.name = "Arial"
        tp.font.size = Pt(18)
        tp.font.bold = True
        tp.font.color.rgb = COLOR_TITLE
        tp.space_after = Pt(4)
        
        stp = ctf.add_paragraph()
        stp.text = subtitle
        stp.font.name = "Arial"
        stp.font.size = Pt(13)
        stp.font.italic = True
        stp.font.color.rgb = COLOR_MUTED
        stp.space_after = Pt(12)
        
        for item in items:
            itp = ctf.add_paragraph()
            itp.text = "▪ " + item
            itp.font.name = "Arial"
            itp.font.size = Pt(13)
            itp.font.color.rgb = COLOR_TEXT
            itp.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 6: Backend Architecture
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "How the Backend Saves and Searches", "05. Backend Architecture")
    
    # 2 Column layout
    col1_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf1 = col1_box.text_frame
    style_text_frame(tf1)
    
    p1 = tf1.paragraphs[0]
    p1.text = "Web Services & Relational DB"
    p1.font.name = "Arial"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TITLE
    p1.space_after = Pt(10)
    
    bullets1 = [
        "Secure Accounts: Encrypts passwords and issues security tokens (JWT) to lock directories.",
        "Trip Storage: Relational SQLite database tables save user profiles and compiled JSON itineraries.",
        "Advisory Scrapers: Automatically parses web pages to download active local weather warnings."
    ]
    for b in bullets1:
        bp = tf1.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)
        
    col2_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf2 = col2_box.text_frame
    style_text_frame(tf2)
    
    p2 = tf2.paragraphs[0]
    p2.text = "RAG Document Retrieval"
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TITLE
    p2.space_after = Pt(10)
    
    bullets2 = [
        "Local Vector Storage: Loads embedded travel paragraphs directly from a JSON file, removing hosting fees.",
        "NumPy Matching: Checks matching coordinates and texts via cosine similarity lookups.",
        "Instruction Injection: Injects matched text directly into AI prompts to prevent fake itineraries."
    ]
    for b in bullets2:
        bp = tf2.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 7: Frontend Architecture
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "The User Interface Design", "06. Frontend Architecture")
    
    # 2 Column layout
    col1_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf1 = col1_box.text_frame
    style_text_frame(tf1)
    
    p1 = tf1.paragraphs[0]
    p1.text = "React SPA & States"
    p1.font.name = "Arial"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TITLE
    p1.space_after = Pt(10)
    
    bullets1 = [
        "Clean Page Switching: Leverages react-router-dom to navigate without annoying page reloads.",
        "Shared Variables: React context distributes authentication tokens and monitors Firebase OTP status.",
        "Interactive Routes: Plots coordinates and draws travel lines dynamically over Leaflet map views."
    ]
    for b in bullets1:
        bp = tf1.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)
        
    col2_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf2 = col2_box.text_frame
    style_text_frame(tf2)
    
    p2 = tf2.paragraphs[0]
    p2.text = "Premium Look and Feel"
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TITLE
    p2.space_after = Pt(10)
    
    bullets2 = [
        "Glassmorphism Styling: Semi-transparent frosted elements to create a premium appearance.",
        "Responsive Scales: Flex grids auto-adjust components to fit desktop screens and mobile interfaces.",
        "Micro-animations: Smooth transitions and hover changes to ensure satisfying user interactions."
    ]
    for b in bullets2:
        bp = tf2.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 8: Libraries Used
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "Key Libraries Powering the App", "07. Libraries Used")
    
    # Two Columns (Python vs JS)
    col1_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf1 = col1_box.text_frame
    style_text_frame(tf1)
    
    p1 = tf1.paragraphs[0]
    p1.text = "Python Server Libraries"
    p1.font.name = "Arial"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TITLE
    p1.space_after = Pt(10)
    
    bullets1 = [
        "fastapi & uvicorn: Coordinates backend web routers and serves APIs.",
        "google-generativeai: Enables native connection to Gemini models.",
        "numpy: Calculates mathematical checks for vector comparisons.",
        "beautifulsoup4: Extracts information from web resources."
    ]
    for b in bullets1:
        bp = tf1.add_paragraph()
        bp.text = "▪ " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)
        
    col2_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf2 = col2_box.text_frame
    style_text_frame(tf2)
    
    p2 = tf2.paragraphs[0]
    p2.text = "JavaScript Frontend Packages"
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TITLE
    p2.space_after = Pt(10)
    
    bullets2 = [
        "react & react-dom: Component base and virtual DOM renderer.",
        "react-router-dom: Moves users cleanly between different page screens.",
        "leaflet & react-leaflet: Visualizes map coordinates and route paths.",
        "canvas-confetti: Confetti animation upon itinerary completion."
    ]
    for b in bullets2:
        bp = tf2.add_paragraph()
        bp.text = "▪ " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 9: Outcomes
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "What We Accomplished", "08. Outcomes")
    
    # 2 Column layout
    col1_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf1 = col1_box.text_frame
    style_text_frame(tf1)
    
    p1 = tf1.paragraphs[0]
    p1.text = "What the App Achieves"
    p1.font.name = "Arial"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TITLE
    p1.space_after = Pt(10)
    
    bullets1 = [
        "Instant Itinerary Synthesis: Generates complete travel dossiers within seconds.",
        "Realistic Budgets: Cost estimators forecast dynamic splits matching budget expectations.",
        "SMS Authentication: SMS validation works smoothly via Firebase OTP onboarding."
    ]
    for b in bullets1:
        bp = tf1.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)
        
    col2_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf2 = col2_box.text_frame
    style_text_frame(tf2)
    
    p2 = tf2.paragraphs[0]
    p2.text = "Zero Hosting Costs"
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TITLE
    p2.space_after = Pt(10)
    
    bullets2 = [
        "Monorepo Setup: Packages the React application directly into FastAPI's static assets directory.",
        "Single-Port Hosting: Avoids complex cross-origin setup issues (CORS).",
        "Minimal Footprint: Relies on minimal CPU limits, allowing hosting on free server tiers."
    ]
    for b in bullets2:
        bp = tf2.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 10: Future Scope & Conclusion
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_slide_header(slide, "Future Scope & Conclusion", "09. Conclusion")
    
    # 2 Column layout
    col1_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf1 = col1_box.text_frame
    style_text_frame(tf1)
    
    p1 = tf1.paragraphs[0]
    p1.text = "Roadmap Ahead"
    p1.font.name = "Arial"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TITLE
    p1.space_after = Pt(10)
    
    bullets1 = [
        "Multi-City Optimizer: Advanced routing algorithms to calculate shortest paths between cities.",
        "Collaborative Portals: Real-time travel rooms so friends can build itineraries together.",
        "Offline Cache Sync: Service workers cache map coordinates for offline travel support."
    ]
    for b in bullets1:
        bp = tf1.add_paragraph()
        bp.text = "• " + b
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)
        
    col2_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf2 = col2_box.text_frame
    style_text_frame(tf2)
    
    p2 = tf2.paragraphs[0]
    p2.text = "Summary"
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TITLE
    p2.space_after = Pt(10)
    
    p_desc = tf2.add_paragraph()
    p_desc.text = "India Tour Planner proves that travel organization doesn't have to be a headache. By bringing together AI generation, relational databases, and interactive maps under a beautiful glassmorphic UI, we have made trip planning simple, intuitive, and fun."
    p_desc.font.name = "Arial"
    p_desc.font.size = Pt(15)
    p_desc.font.color.rgb = COLOR_TEXT
    
    # Save to file
    out_filename = "India_Tour_Planner_Presentation.pptx"
    out_path = os.path.join("c:\\Users\\Acer\\OneDrive\\Desktop\\India Tour", out_filename)
    prs.save(out_path)
    print(f"\n[+] Success! Presentation generated at:\n    {out_path}\n")

if __name__ == "__main__":
    create_presentation()
